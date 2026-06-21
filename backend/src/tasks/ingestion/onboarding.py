"""
Onboarding Materials Ingestion Celery Tasks

This module implements Celery tasks for ingesting onboarding materials into the knowledge base.
Handles extracting text from PDFs, Word documents, PowerPoint presentations, and Markdown files
stored in MinIO, chunking content, generating embeddings, and database storage.
"""

import io
import logging
import hashlib
from datetime import datetime
from typing import Dict, Any, List, Optional
import asyncio

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

# Document processing libraries
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
import markdown

from tasks.celery import celery_app
from core.database import get_session_factory
from core.config import settings
from core.minio_client import download_file, get_file_info, BUCKET_KNOWLEDGE_FILES
from models.data_source import DataSource, IngestionJob, JobStatus, DataSourceType
from models.knowledge import KnowledgeDocument, DocumentEmbedding, ContentType
from integrations.llm_factory import LLMFactory
from integrations.openai_client import OpenAIError
from integrations.ollama_client import OllamaError
from utils.text_processing import chunk_text, count_tokens, clean_html

# Logger
logger = logging.getLogger(__name__)

# Constants
CHUNK_SIZE = 500  # tokens
CHUNK_OVERLAP = 50  # tokens
EMBEDDING_BATCH_SIZE = 100  # Process embeddings in batches
SUPPORTED_FORMATS = ['.pdf', '.docx', '.pptx', '.md']


# ============================================================================
# Celery Tasks
# ============================================================================

@celery_app.task(
    bind=True,
    name="tasks.ingestion.onboarding.ingest_onboarding_materials",
    max_retries=3,
    default_retry_delay=60,
)
def ingest_onboarding_materials(self, data_source_id: int, file_path: str) -> Dict[str, Any]:
    """
    Ingest onboarding materials from MinIO storage.
    
    This task:
    1. Retrieves file from MinIO knowledge-files bucket
    2. Extracts text based on file type (PDF/Word/PowerPoint/Markdown)
    3. Chunks documents into 500-token chunks with 50-token overlap
    4. Generates embeddings for each chunk (batch processing)
    5. Stores documents in knowledge_documents table with metadata
    6. Stores embeddings in document_embeddings table
    7. Updates ingestion_jobs table with progress and status
    8. Performs deduplication using document_hash (SHA256)
    
    Args:
        data_source_id: ID of the onboarding data source
        file_path: Path to file in MinIO bucket (e.g., "onboarding/guidelines.pdf")
        
    Returns:
        dict: Summary with status, documents processed, documents failed
        
    Raises:
        Exception: Retries on failure with exponential backoff
    """
    try:
        # Reset global database engine to avoid event loop conflicts in Celery workers
        import core.database as db_module
        db_module._engine = None
        db_module._async_session_factory = None

        # Run async ingestion in event loop
        # Use asyncio.run() instead of get_event_loop() to avoid event loop conflicts in Celery
        result = asyncio.run(
            _ingest_onboarding_materials_async(data_source_id, file_path, self)
        )
        return result
    except Exception as exc:
        logger.error(f"Onboarding materials ingestion failed: {exc}", exc_info=True)
        # Retry with exponential backoff
        raise self.retry(exc=exc, countdown=2 ** self.request.retries * 60)
    finally:
        # Clean up the engine created in this async context
        import core.database as db_module
        db_module._engine = None
        db_module._async_session_factory = None


# ============================================================================
# Async Implementation
# ============================================================================

async def _ingest_onboarding_materials_async(
    data_source_id: int,
    file_path: str,
    task_context
) -> Dict[str, Any]:
    """
    Async implementation of onboarding materials ingestion.
    
    Args:
        data_source_id: ID of the onboarding data source
        file_path: Path to file in MinIO bucket
        task_context: Celery task context for updates
        
    Returns:
        dict: Summary with status, job_id, documents_processed, documents_failed
    """
    session_factory = get_session_factory()
    
    async with session_factory() as session:
        try:
            # Fetch data source
            data_source = await _get_data_source(session, data_source_id)
            
            if not data_source:
                raise ValueError(f"Data source {data_source_id} not found")
                
            if data_source.type != DataSourceType.ONBOARDING:
                raise ValueError(f"Data source {data_source_id} is not an onboarding source")
                
            if not data_source.is_active:
                logger.warning(f"Data source {data_source_id} is not active, skipping")
                return {"status": "skipped", "documents_processed": 0, "documents_failed": 0}
            
            # Create ingestion job
            job = await _create_ingestion_job(session, data_source_id)
            job_id = job.id
            await session.commit()
            
            logger.info(f"Starting onboarding materials ingestion for file: {file_path}")
            
            # Process the file
            try:
                await _process_file(session, data_source, file_path, job_id)

                # Mark job as successful
                await _complete_job(session, job_id, JobStatus.SUCCESS, 1, 0)
                await session.commit()

                logger.info(f"Onboarding materials ingestion completed: {file_path}")

                return {
                    "status": "success",
                    "job_id": job_id,
                    "documents_processed": 1,
                    "documents_failed": 0,
                    "file_path": file_path
                }

            except Exception as e:
                logger.error(f"Failed to process file {file_path}: {e}", exc_info=True)
                await _complete_job(session, job_id, JobStatus.FAILED, 0, 1, str(e))
                await session.commit()

                return {
                    "status": "failed",
                    "job_id": job_id,
                    "documents_processed": 0,
                    "documents_failed": 1,
                    "error": str(e)
                }

        except Exception as e:
            logger.error(f"Onboarding materials ingestion error: {e}", exc_info=True)
            await session.rollback()
            raise


# ============================================================================
# Helper Functions
# ============================================================================

async def _get_data_source(session: AsyncSession, data_source_id: int) -> Optional[DataSource]:
    """Fetch data source by ID."""
    result = await session.execute(
        select(DataSource).where(DataSource.id == data_source_id)
    )
    return result.scalar_one_or_none()


async def _create_ingestion_job(session: AsyncSession, data_source_id: int) -> IngestionJob:
    """Create a new ingestion job record."""
    job = IngestionJob(
        data_source_id=data_source_id,
        status=JobStatus.RUNNING,
        documents_processed=0,
        documents_failed=0,
        started_at=datetime.utcnow()
    )
    session.add(job)
    await session.flush()
    return job


async def _complete_job(
    session: AsyncSession,
    job_id: int,
    status: JobStatus,
    processed: int,
    failed: int,
    error_message: Optional[str] = None
):
    """Mark ingestion job as complete."""
    result = await session.execute(
        select(IngestionJob).where(IngestionJob.id == job_id)
    )
    job = result.scalar_one_or_none()

    if job:
        job.status = status
        job.documents_processed = processed
        job.documents_failed = failed
        job.completed_at = datetime.utcnow()
        if error_message:
            job.error_message = error_message[:500]  # Truncate long errors


async def _process_file(
    session: AsyncSession,
    data_source: DataSource,
    file_path: str,
    job_id: int
):
    """
    Process a single file: extract text, chunk, embed, and store.

    Args:
        session: Database session
        data_source: Data source configuration
        file_path: Path to file in MinIO
        job_id: Ingestion job ID
    """
    # Get file metadata from MinIO
    file_info = get_file_info(BUCKET_KNOWLEDGE_FILES, file_path)
    if not file_info:
        raise ValueError(f"File not found in MinIO: {file_path}")

    logger.info(f"Processing file: {file_path} ({file_info['size']} bytes)")

    # Download file from MinIO
    file_stream = download_file(BUCKET_KNOWLEDGE_FILES, file_path)
    if not file_stream:
        raise ValueError(f"Failed to download file: {file_path}")

    # Read file content
    file_content = file_stream.read()
    file_stream.close()

    # Extract text based on file extension
    file_extension = _get_file_extension(file_path)

    if file_extension == '.pdf':
        text = _extract_text_from_pdf(file_content)
    elif file_extension == '.docx':
        text = _extract_text_from_docx(file_content)
    elif file_extension == '.pptx':
        text = _extract_text_from_pptx(file_content)
    elif file_extension == '.md':
        text = _extract_text_from_markdown(file_content)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

    if not text or not text.strip():
        raise ValueError(f"No text extracted from file: {file_path}")

    logger.info(f"Extracted {len(text)} characters from {file_path}")

    # Compute document hash for deduplication
    document_hash = _compute_document_hash(text)

    # Check for existing document with same hash
    # Note: Query the database directly to check for duplicates
    data_source_id = data_source.id

    # Query for existing document with same data_source_id and hash
    # Use execution_options to bypass session cache and query database directly
    result = await session.execute(
        select(KnowledgeDocument).where(
            KnowledgeDocument.data_source_id == data_source_id,
            KnowledgeDocument.document_hash == document_hash,
            KnowledgeDocument.is_deleted == False
        ).execution_options(populate_existing=True)
    )
    existing = result.scalar_one_or_none()

    if existing:
        logger.info(f"Document already exists (hash match), skipping: {file_path}")
        return

    # Create knowledge document
    filename = file_path.split('/')[-1]
    document = KnowledgeDocument(
        data_source_id=data_source_id,  # Use captured ID to avoid lazy-load
        external_id=file_path,  # Use file path as external ID
        title=filename,
        content=text,
        content_type=ContentType.DOCUMENT,
        url=None,  # MinIO files don't have public URLs
        doc_metadata={
            "filename": filename,
            "file_path": file_path,
            "file_size": file_info['size'],
            "content_type": file_info['content_type'],
            "file_extension": file_extension,
            "last_modified": file_info['last_modified'].isoformat() if file_info['last_modified'] else None,
        },
        document_hash=document_hash,
        indexed_at=datetime.utcnow()
    )

    session.add(document)
    await session.flush()  # Get document ID

    logger.info(f"Created knowledge document: {document.id}")

    # Chunk the text
    chunks = chunk_text(text, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP)
    logger.info(f"Split document into {len(chunks)} chunks")

    # Generate embeddings for chunks
    llm_client = LLMFactory.create_client()

    try:
        # Check if client supports batch_size parameter (OpenAI does, Ollama doesn't)
        if hasattr(llm_client, '__class__') and llm_client.__class__.__name__ == 'OpenAIClient':
            embeddings = await llm_client.generate_embeddings_batch(
                chunks,
                batch_size=EMBEDDING_BATCH_SIZE
            )
        else:
            embeddings = await llm_client.generate_embeddings_batch(chunks)

        logger.info(f"Generated {len(embeddings)} embeddings")

        # Store embeddings in database
        for idx, (chunk_text_content, embedding) in enumerate(zip(chunks, embeddings)):
            token_count = count_tokens(chunk_text_content)

            embedding_record = DocumentEmbedding(
                document_id=document.id,
                chunk_index=idx,
                chunk_text=chunk_text_content,
                embedding=embedding,
                token_count=token_count
            )
            session.add(embedding_record)

        logger.info(f"Stored {len(embeddings)} embeddings for document {document.id}")

    except (OpenAIError, OllamaError) as e:
        logger.error(f"Failed to generate embeddings: {e}")
        raise


def _get_file_extension(file_path: str) -> str:
    """Extract file extension from path."""
    if '.' not in file_path:
        return ''
    return '.' + file_path.split('.')[-1].lower()


def _extract_text_from_pdf(file_content: bytes) -> str:
    """
    Extract text from PDF file using PyPDF2.

    Args:
        file_content: PDF file bytes

    Returns:
        str: Extracted text
    """
    try:
        pdf_file = io.BytesIO(file_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        text_parts = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            if text:
                text_parts.append(text)

        full_text = '\n\n'.join(text_parts)
        logger.debug(f"Extracted {len(full_text)} characters from PDF ({len(pdf_reader.pages)} pages)")

        return full_text

    except Exception as e:
        logger.error(f"Failed to extract text from PDF: {e}")
        raise ValueError(f"PDF extraction failed: {e}")


def _extract_text_from_docx(file_content: bytes) -> str:
    """
    Extract text from Word document using python-docx.

    Args:
        file_content: DOCX file bytes

    Returns:
        str: Extracted text
    """
    try:
        docx_file = io.BytesIO(file_content)
        doc = DocxDocument(docx_file)

        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)

        full_text = '\n\n'.join(text_parts)
        logger.debug(f"Extracted {len(full_text)} characters from DOCX ({len(doc.paragraphs)} paragraphs)")

        return full_text

    except Exception as e:
        logger.error(f"Failed to extract text from DOCX: {e}")
        raise ValueError(f"DOCX extraction failed: {e}")


def _extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extract text from PowerPoint presentation using python-pptx.

    Args:
        file_content: PPTX file bytes

    Returns:
        str: Extracted text
    """
    try:
        pptx_file = io.BytesIO(file_content)
        prs = Presentation(pptx_file)

        text_parts = []
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                text_parts.append(f"[Slide {slide_num + 1}]\n" + '\n'.join(slide_text))

        full_text = '\n\n'.join(text_parts)
        logger.debug(f"Extracted {len(full_text)} characters from PPTX ({len(prs.slides)} slides)")

        return full_text

    except Exception as e:
        logger.error(f"Failed to extract text from PPTX: {e}")
        raise ValueError(f"PPTX extraction failed: {e}")


def _extract_text_from_markdown(file_content: bytes) -> str:
    """
    Extract text from Markdown file.

    Converts Markdown to HTML, then extracts plain text.

    Args:
        file_content: Markdown file bytes

    Returns:
        str: Extracted text
    """
    try:
        # Decode bytes to string
        markdown_text = file_content.decode('utf-8')

        # Convert Markdown to HTML
        html = markdown.markdown(markdown_text)

        # Extract plain text from HTML
        plain_text = clean_html(html)

        logger.debug(f"Extracted {len(plain_text)} characters from Markdown")

        return plain_text

    except Exception as e:
        logger.error(f"Failed to extract text from Markdown: {e}")
        raise ValueError(f"Markdown extraction failed: {e}")


def _compute_document_hash(content: str) -> str:
    """
    Compute SHA-256 hash of document content for deduplication.

    Args:
        content: Document text content

    Returns:
        str: Hexadecimal hash string
    """
    return hashlib.sha256(content.encode('utf-8')).hexdigest()
