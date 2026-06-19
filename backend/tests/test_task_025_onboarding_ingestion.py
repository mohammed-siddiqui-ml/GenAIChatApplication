"""
Tests for Onboarding Materials Ingestion (Task 025)

This test suite covers:
- Unit tests for helper functions (text extraction, hashing, file extension)
- Integration tests for full pipeline (PDF, DOCX, PPTX, Markdown)
- Functional tests for business logic (deduplication, job tracking)
- Error handling tests (missing files, invalid formats, API failures)
"""

import io
import pytest
import pytest_asyncio
from unittest.mock import Mock, AsyncMock, patch, MagicMock
from datetime import datetime

# PDF, DOCX, PPTX creation libraries for test data
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches, Pt

from models.data_source import DataSource, IngestionJob, JobStatus, DataSourceType
from models.knowledge import KnowledgeDocument, DocumentEmbedding, ContentType
from tasks.ingestion.onboarding import (
    _get_file_extension,
    _compute_document_hash,
    _extract_text_from_pdf,
    _extract_text_from_docx,
    _extract_text_from_pptx,
    _extract_text_from_markdown,
    _get_data_source,
    _create_ingestion_job,
    _complete_job,
    _process_file,
    _ingest_onboarding_materials_async,
)


# ============================================================================
# Test Data Generators
# ============================================================================

def create_test_pdf(text: str) -> bytes:
    """Create a simple PDF with given text."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    c.drawString(100, 750, text)
    c.save()
    buffer.seek(0)
    return buffer.read()


def create_test_docx(paragraphs: list) -> bytes:
    """Create a DOCX with given paragraphs."""
    doc = DocxDocument()
    for para in paragraphs:
        doc.add_paragraph(para)
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def create_test_pptx(slides_text: list) -> bytes:
    """Create a PPTX with given slides (list of lists of text)."""
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    
    for slide_content in slides_text:
        slide = prs.slides.add_slide(blank_slide_layout)
        for idx, text in enumerate(slide_content):
            txBox = slide.shapes.add_textbox(
                Inches(1), Inches(1 + idx * 0.5), Inches(8), Inches(0.4)
            )
            tf = txBox.text_frame
            tf.text = text
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def create_test_markdown(content: str) -> bytes:
    """Create Markdown file bytes."""
    return content.encode('utf-8')


# ============================================================================
# Unit Tests - Helper Functions
# ============================================================================

@pytest.mark.unit
class TestUnitFunctions:
    """Test individual helper functions in isolation."""
    
    def test_get_file_extension(self):
        """Test file extension extraction (TC-025-U01)."""
        assert _get_file_extension("onboarding/guide.pdf") == ".pdf"
        assert _get_file_extension("folder/subfolder/document.docx") == ".docx"
        assert _get_file_extension("presentation.PPTX") == ".pptx"
        assert _get_file_extension("readme.MD") == ".md"
        assert _get_file_extension("archive.tar.gz") == ".gz"
        assert _get_file_extension("no_extension") == ""
    
    def test_compute_document_hash(self):
        """Test document hash computation (TC-025-U02)."""
        hash1 = _compute_document_hash("Hello World")
        hash2 = _compute_document_hash("Hello World")
        hash3 = _compute_document_hash("Goodbye World")
        
        # Same content produces same hash
        assert hash1 == hash2
        # Different content produces different hash
        assert hash1 != hash3
        # Hash format: 64-character hexadecimal
        assert len(hash1) == 64
        assert all(c in '0123456789abcdef' for c in hash1)
    
    def test_extract_text_from_pdf_success(self):
        """Test PDF text extraction (TC-025-U03)."""
        pdf_bytes = create_test_pdf("Sample PDF Content")
        text = _extract_text_from_pdf(pdf_bytes)
        
        assert "Sample PDF Content" in text
        assert len(text) > 0
    
    def test_extract_text_from_pdf_corrupted(self):
        """Test PDF extraction with corrupted file (TC-025-U05)."""
        with pytest.raises(ValueError, match="PDF extraction failed"):
            _extract_text_from_pdf(b"not a pdf")
    
    def test_extract_text_from_docx_success(self):
        """Test DOCX text extraction (TC-025-U06)."""
        docx_bytes = create_test_docx(["Introduction", "Body", "Conclusion"])
        text = _extract_text_from_docx(docx_bytes)

        assert "Introduction" in text
        assert "Body" in text
        assert "Conclusion" in text

    def test_extract_text_from_docx_corrupted(self):
        """Test DOCX extraction with corrupted file."""
        with pytest.raises(ValueError, match="DOCX extraction failed"):
            _extract_text_from_docx(b"not a docx")

    def test_extract_text_from_pptx_success(self):
        """Test PPTX text extraction (TC-025-U07)."""
        pptx_bytes = create_test_pptx([["Welcome", "Introduction"], ["Summary", "Thank You"]])
        text = _extract_text_from_pptx(pptx_bytes)

        assert "[Slide 1]" in text
        assert "[Slide 2]" in text
        assert "Welcome" in text
        assert "Thank You" in text

    def test_extract_text_from_pptx_corrupted(self):
        """Test PPTX extraction with corrupted file."""
        with pytest.raises(ValueError, match="PPTX extraction failed"):
            _extract_text_from_pptx(b"not a pptx")

    def test_extract_text_from_markdown_success(self):
        """Test Markdown text extraction (TC-025-U08)."""
        md_content = "# Header\n\n**Bold** text\n\n- List item"
        md_bytes = create_test_markdown(md_content)
        text = _extract_text_from_markdown(md_bytes)

        # Should have plain text without Markdown formatting
        assert "Header" in text
        assert "Bold" in text or "text" in text
        assert "List item" in text

    def test_extract_text_from_markdown_invalid_encoding(self):
        """Test Markdown extraction with invalid encoding."""
        # Invalid UTF-8 bytes
        invalid_bytes = b'\xff\xfe\xfd'
        with pytest.raises(ValueError, match="Markdown extraction failed"):
            _extract_text_from_markdown(invalid_bytes)


# ============================================================================
# Unit Tests - Database Helpers
# ============================================================================

@pytest.mark.unit
@pytest.mark.asyncio
class TestDatabaseHelpers:
    """Test database helper functions."""

    async def test_get_data_source_success(self, db_session):
        """Test successful data source retrieval (TC-025-U09)."""
        # Create test data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Retrieve it
        result = await _get_data_source(db_session, data_source.id)

        assert result is not None
        assert result.type == DataSourceType.ONBOARDING
        assert result.is_active is True

    async def test_get_data_source_not_found(self, db_session):
        """Test data source not found (TC-025-U10)."""
        result = await _get_data_source(db_session, 99999)
        assert result is None

    async def test_create_ingestion_job(self, db_session):
        """Test ingestion job creation (TC-025-U11)."""
        # Create test data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create job
        job = await _create_ingestion_job(db_session, data_source.id)
        await db_session.commit()
        await db_session.refresh(job)

        assert job.status == JobStatus.RUNNING
        assert job.started_at is not None
        assert job.documents_processed == 0
        assert job.documents_failed == 0
        assert job.completed_at is None

    async def test_complete_job_success(self, db_session):
        """Test job completion with success (TC-025-U12)."""
        # Create data source and job
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={}
        )
        db_session.add(data_source)
        await db_session.commit()

        job = await _create_ingestion_job(db_session, data_source.id)
        await db_session.commit()
        await db_session.refresh(job)

        # Complete job
        await _complete_job(db_session, job.id, JobStatus.SUCCESS, 5, 0)
        await db_session.commit()
        await db_session.refresh(job)

        assert job.status == JobStatus.SUCCESS
        assert job.documents_processed == 5
        assert job.documents_failed == 0
        assert job.completed_at is not None
        assert job.error_message is None

    async def test_complete_job_failure(self, db_session):
        """Test job completion with failure (TC-025-U13)."""
        # Create data source and job
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={}
        )
        db_session.add(data_source)
        await db_session.commit()

        job = await _create_ingestion_job(db_session, data_source.id)
        await db_session.commit()
        await db_session.refresh(job)

        # Complete job with error
        await _complete_job(db_session, job.id, JobStatus.FAILED, 2, 3, "OpenAI API error")
        await db_session.commit()
        await db_session.refresh(job)

        assert job.status == JobStatus.FAILED
        assert job.documents_processed == 2
        assert job.documents_failed == 3
        assert job.error_message == "OpenAI API error"


# ============================================================================
# Integration Tests - Full Pipeline
# ============================================================================

@pytest.mark.integration
@pytest.mark.asyncio
class TestFullPipeline:
    """Test end-to-end file processing pipeline."""

    @pytest.fixture
    def mock_file_info(self):
        """Mock MinIO file info."""
        return {
            "bucket": "knowledge-files",
            "object_name": "onboarding/guide.pdf",
            "size": 12345,
            "etag": "abc123",
            "content_type": "application/pdf",
            "last_modified": datetime(2024, 1, 15, 10, 0, 0),
            "metadata": {}
        }

    @pytest.fixture
    def mock_embedding(self):
        """Mock OpenAI embedding vector."""
        return [0.1] * 1536  # 1536-dimensional vector

    async def test_pdf_ingestion(self, db_session, mock_file_info, mock_embedding):
        """Test full PDF ingestion pipeline (TC-025-I01)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create PDF content
        pdf_bytes = create_test_pdf("This is a test PDF document for onboarding materials.")

        # Mock MinIO and OpenAI
        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download, \
             patch('tasks.ingestion.onboarding.OpenAIClient') as mock_openai_class:

            # Setup mocks
            mock_get_info.return_value = mock_file_info
            mock_stream = io.BytesIO(pdf_bytes)
            mock_download.return_value = mock_stream

            # Mock OpenAI client
            mock_openai = AsyncMock()
            mock_openai.generate_embeddings_batch.return_value = [mock_embedding]
            mock_openai_class.return_value = mock_openai

            # Process file
            await _process_file(db_session, data_source, "onboarding/guide.pdf", 1)
            await db_session.commit()

            # Verify document created
            from sqlalchemy import select
            result = await db_session.execute(select(KnowledgeDocument))
            documents = result.scalars().all()

            assert len(documents) == 1
            doc = documents[0]
            assert doc.external_id == "onboarding/guide.pdf"
            assert doc.title == "guide.pdf"
            assert doc.content_type == ContentType.DOCUMENT
            assert "test PDF document" in doc.content
            assert doc.document_hash is not None
            assert doc.doc_metadata["filename"] == "guide.pdf"
            assert doc.doc_metadata["file_size"] == 12345

            # Verify embeddings created
            result = await db_session.execute(select(DocumentEmbedding))
            embeddings = result.scalars().all()

            assert len(embeddings) >= 1
            emb = embeddings[0]
            assert emb.document_id == doc.id
            assert emb.chunk_index == 0
            assert emb.embedding == mock_embedding
            assert emb.token_count > 0

    async def test_docx_ingestion(self, db_session, mock_file_info, mock_embedding):
        """Test full DOCX ingestion pipeline (TC-025-I02)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create DOCX content
        docx_bytes = create_test_docx(["Employee Handbook", "Welcome to our company!", "Code of Conduct"])

        # Update file info for DOCX
        file_info = mock_file_info.copy()
        file_info["content_type"] = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

        # Mock MinIO and OpenAI
        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download, \
             patch('tasks.ingestion.onboarding.OpenAIClient') as mock_openai_class:

            # Setup mocks
            mock_get_info.return_value = file_info
            mock_stream = io.BytesIO(docx_bytes)
            mock_download.return_value = mock_stream

            # Mock OpenAI client
            mock_openai = AsyncMock()
            mock_openai.generate_embeddings_batch.return_value = [mock_embedding]
            mock_openai_class.return_value = mock_openai

            # Process file
            await _process_file(db_session, data_source, "onboarding/handbook.docx", 1)
            await db_session.commit()

            # Verify document created
            from sqlalchemy import select
            result = await db_session.execute(select(KnowledgeDocument))
            documents = result.scalars().all()

            assert len(documents) == 1
            doc = documents[0]
            assert doc.external_id == "onboarding/handbook.docx"
            assert "Employee Handbook" in doc.content
            assert "Welcome" in doc.content

    async def test_deduplication_same_hash(self, db_session, mock_file_info, mock_embedding):
        """Test deduplication with same document hash (TC-025-F01)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create PDF content
        pdf_bytes = create_test_pdf("Same content for deduplication test")

        # Mock MinIO and OpenAI
        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download, \
             patch('tasks.ingestion.onboarding.OpenAIClient') as mock_openai_class:

            # Setup mocks
            mock_get_info.return_value = mock_file_info
            mock_stream = io.BytesIO(pdf_bytes)
            mock_download.return_value = mock_stream

            # Mock OpenAI client
            mock_openai = AsyncMock()
            mock_openai.generate_embeddings_batch.return_value = [mock_embedding]
            mock_openai_class.return_value = mock_openai

            # Process file first time
            await _process_file(db_session, data_source, "onboarding/guide.pdf", 1)
            await db_session.commit()

            # CRITICAL: Refresh the data_source object after commit to ensure it's up-to-date
            await db_session.refresh(data_source)

            # Reset stream for second processing
            mock_stream = io.BytesIO(pdf_bytes)
            mock_download.return_value = mock_stream

            # Process same file again (should skip due to deduplication)
            await _process_file(db_session, data_source, "onboarding/guide.pdf", 1)
            await db_session.commit()

            # Verify only one document exists
            from sqlalchemy import select
            result = await db_session.execute(select(KnowledgeDocument))
            documents = result.scalars().all()

            assert len(documents) == 1  # Only one document should exist


# ============================================================================
# Error Handling Tests
# ============================================================================

@pytest.mark.error_handling
@pytest.mark.asyncio
class TestErrorHandling:
    """Test error handling and failure scenarios."""

    async def test_missing_file_in_minio(self, db_session):
        """Test handling of missing file in MinIO (TC-025-E01)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Mock MinIO to return None (file not found)
        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info:
            mock_get_info.return_value = None

            # Should raise ValueError
            with pytest.raises(ValueError, match="File not found in MinIO"):
                await _process_file(db_session, data_source, "onboarding/missing.pdf", 1)

    async def test_unsupported_file_format(self, db_session):
        """Test handling of unsupported file format (TC-025-E02)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Mock MinIO with .txt file
        file_info = {
            "bucket": "knowledge-files",
            "object_name": "onboarding/document.txt",
            "size": 100,
            "content_type": "text/plain",
            "last_modified": datetime.utcnow(),
        }

        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download:

            mock_get_info.return_value = file_info
            mock_stream = io.BytesIO(b"Plain text content")
            mock_download.return_value = mock_stream

            # Should raise ValueError for unsupported format
            with pytest.raises(ValueError, match="Unsupported file format: .txt"):
                await _process_file(db_session, data_source, "onboarding/document.txt", 1)

    async def test_empty_file_no_text(self, db_session):
        """Test handling of empty file with no extractable text (TC-025-E03)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create empty PDF
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        buffer = io.BytesIO()
        c = canvas.Canvas(buffer, pagesize=letter)
        c.save()  # Save without adding any text
        buffer.seek(0)
        empty_pdf = buffer.read()

        file_info = {
            "bucket": "knowledge-files",
            "object_name": "onboarding/empty.pdf",
            "size": len(empty_pdf),
            "content_type": "application/pdf",
            "last_modified": datetime.utcnow(),
        }

        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download:

            mock_get_info.return_value = file_info
            mock_stream = io.BytesIO(empty_pdf)
            mock_download.return_value = mock_stream

            # Should raise ValueError for no text
            with pytest.raises(ValueError, match="No text extracted from file"):
                await _process_file(db_session, data_source, "onboarding/empty.pdf", 1)


# ============================================================================
# Functional Tests - End-to-End
# ============================================================================

@pytest.mark.functional
@pytest.mark.asyncio
class TestFunctionalWorkflows:
    """Test complete business workflows end-to-end."""

    async def test_complete_ingestion_workflow(self, db_session):
        """Test complete ingestion workflow with job tracking (TC-025-F04)."""
        # Create data source
        data_source = DataSource(
            name="Test Onboarding",
            type=DataSourceType.ONBOARDING,
            is_active=True,
            source_config={"bucket": "knowledge-files"}
        )
        db_session.add(data_source)
        await db_session.commit()
        await db_session.refresh(data_source)

        # Create PDF
        pdf_bytes = create_test_pdf("Complete workflow test document")

        file_info = {
            "bucket": "knowledge-files",
            "object_name": "onboarding/workflow.pdf",
            "size": len(pdf_bytes),
            "content_type": "application/pdf",
            "last_modified": datetime.utcnow(),
        }

        mock_embedding = [0.1] * 1536

        # Create a mock session factory that returns the test session
        from contextlib import asynccontextmanager

        @asynccontextmanager
        async def mock_session_factory():
            """Mock session factory that yields the test session."""
            yield db_session

        # Mock MinIO, OpenAI, and get_session_factory
        with patch('tasks.ingestion.onboarding.get_file_info') as mock_get_info, \
             patch('tasks.ingestion.onboarding.download_file') as mock_download, \
             patch('tasks.ingestion.onboarding.OpenAIClient') as mock_openai_class, \
             patch('tasks.ingestion.onboarding.get_session_factory') as mock_get_session_factory:

            # Setup mocks
            mock_get_info.return_value = file_info
            mock_stream = io.BytesIO(pdf_bytes)
            mock_download.return_value = mock_stream

            # Mock OpenAI client
            mock_openai = AsyncMock()
            mock_openai.generate_embeddings_batch.return_value = [mock_embedding]
            mock_openai_class.return_value = mock_openai

            # Mock session factory to return test session
            mock_get_session_factory.return_value = mock_session_factory

            # Run complete workflow
            result = await _ingest_onboarding_materials_async(
                data_source.id,
                "onboarding/workflow.pdf",
                None
            )

            # Verify result
            assert result["status"] == "success"
            assert result["documents_processed"] == 1
            assert result["documents_failed"] == 0
            assert "job_id" in result

            # Verify job was created and completed
            from sqlalchemy import select
            job_result = await db_session.execute(
                select(IngestionJob).where(IngestionJob.id == result["job_id"])
            )
            job = job_result.scalar_one()

            assert job.status == JobStatus.SUCCESS
            assert job.documents_processed == 1
            assert job.documents_failed == 0
            assert job.started_at is not None
            assert job.completed_at is not None
